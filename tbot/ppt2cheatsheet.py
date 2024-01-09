# transfer all pptx into cheatsheets.
# argparse,python-pptx,python-docx
import os
import argparse
import sys
import re

from pptx import Presentation
from docx import Document


class Slides2CheatSheet():
    def __init__(self, slides_dir, slides_num=0):
        self.slides = None
        self.slides_text = None
        self.slides_dir = slides_dir
        self.slides_num = slides_num
        print("There are" + str(self.slides_num) + "different slides under dir" + self.slides_dir + ".")
        return

    def read_slides(self):
        # TODO：read and add name of each slides and headline of each slide?
        if len(self.slides) > 0:
            print("start to transfer pptx to docx")
            text_runs = []
            for slide in self.slides:
                prs = Presentation(slide)
                pattern = r'/([^/]+)\.pptx$'
                match = re.search(pattern, slide)
                if match:
                    # Extract information from the matched group
                    extracted_text = match.group(1)
                    text_runs.append(extracted_text+":  ")
                for slide in prs.slides:
                    text_runs.append("[")
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs.append(run.text.replace("\n", " "))
                    text_runs.append("]\n")
            self.slides_text = text_runs
        return

    def transfer(self):
        slides = []
        for filename in os.listdir(self.slides_dir):
            slide = os.path.join(self.slides_dir, filename)
            extension = filename.split(".")[-1]
            if extension == "pptx":
                slides.append(slide)
            else:
                sys.exit("dir contains file " + slide + " that's not ppt or pptx... please remove them")
        self.slides = slides
        print("find" + str(len(slides)) + "slides under this dir...")
        return slides

    def purify(self):
        # TODO: work on how to purify text in general
        return

    def save(self):
        # TODO: should change general space for a letter to occupy whole page
        if self.slides_text is None:
            pass
        else:
            print("Transfering to docx...")
            document = Document()
            p = document.add_paragraph('Start:')
            for i in self.slides_text:
                p.add_run(i)
            for section in document.sections:
                header = section.header
                footer = section.footer
                header.is_linked_to_previous = True
                footer.is_linked_to_previous = True

                from docx.shared import Inches
                section.page_width, section.page_height = Inches(8.5), Inches(11)
                section.left_margin = Inches(0.1)
                section.right_margin = Inches(0.1)
                section.top_margin, section.bottom_margin = 200,200
            document.save('demo.docx')
        print("finished!")
        return


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='This program is designed to extract all slides '
                                                 'and put it into one double paged letter sized docx file.')
    test_dir = "./test"
    files = Slides2CheatSheet(test_dir)
    files.transfer()
    files.read_slides()
    files.save()